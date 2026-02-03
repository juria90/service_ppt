"""PowerPoint automation using python-pptx library.

This module implements App and Presentation classes for manipulating PowerPoint
files using the python-pptx library, providing cross-platform support without
requiring PowerPoint to be installed. Note that image export features are not
available with this implementation.
"""

import copy
import re
from collections.abc import Callable
from enum import IntEnum
from io import BytesIO
from typing import Any

from pptx import Presentation as PptxPres
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE
from pptx.parts.chart import ChartPart
from pptx.parts.embeddedpackage import EmbeddedXlsxPart

from service_ppt.powerpoint_base import PPTAppBase, PresentationBase, SlideCache


class SlideLayout(IntEnum):
    SLD_LAYOUT_TITLE = 0
    SLD_LAYOUT_TITLE_AND_CONTENT = 1
    SLD_LAYOUT_SECTION_HEADER = 2
    SLD_LAYOUT_TWO_CONTENT = 3
    SLD_LAYOUT_COMPARISON = 4
    SLD_LAYOUT_TITLE_ONLY = 5
    SLD_LAYOUT_BLANK = 6
    SLD_LAYOUT_CONTENT_WITH_CAPTION = 7
    SLD_LAYOUT_PICTURE_WITH_CAPTION = 8


def iterate_text_objects_in_shapes(shapes, fn: Callable | None = None):
    for shape in shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if fn:
                        try:
                            result = fn(run)
                            if result:
                                yield run
                        except StopIteration:
                            return
                    else:
                        yield run

        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if fn:
                        try:
                            result = fn(cell)
                            if result:
                                yield cell
                        except StopIteration:
                            return
                    else:
                        yield cell

        if shape.has_chart:
            chart = shape.chart
            if chart.has_title:
                chart_title = chart.chart_title
                if chart_title.has_text_frame:
                    text_frame = chart_title.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            if fn:
                                try:
                                    result = fn(run)
                                    if result:
                                        yield run
                                except StopIteration:
                                    return
                            else:
                                yield run


def get_texts_in_shapes(shapes):
    texts = [obj.text for obj in iterate_text_objects_in_shapes(shapes)]
    return texts


def find_text_in_shapes(shapes, text, match_case=True, whole_words=True):
    # https://stackoverflow.com/questions/25483114/regex-to-find-whole-word-in-text-but-case-insensitive
    pattern = ""
    if match_case is False:
        pattern += r"(?i)"
    if whole_words:
        pattern += r"\b" + re.escape(text) + r"\b"
    else:
        pattern += re.escape(text)

    re_text = re.compile(pattern)

    def match_string(obj):
        nonlocal re_text
        return re_text.search(obj.text) is not None

    texts = [obj.text for obj in iterate_text_objects_in_shapes(shapes, match_string)]
    for shape in shapes:
        if shape.has_text_frame:
            found_text = text in shape.text
            if found_text is not None:
                return found_text

    return texts[0] if len(texts) > 0 else None


def replace_texts_in_shapes(shapes, text_dict):
    count = 0

    def replace_texts_in_run(obj):
        nonlocal count, text_dict
        for from_text, to_text in text_dict.items():
            if (obj.text.find(from_text)) != -1:
                obj.text = obj.text.replace(from_text, to_text)
                count = count + 1

    for _ in iterate_text_objects_in_shapes(shapes, replace_texts_in_run):
        pass

    return count


def dump_shapes(shapes):
    for i, shape in enumerate(shapes):
        if shape.has_text_frame:
            print("shape %d: %d" % (i, shape.shape_type))

            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                print(f"PlaceholderFormat.Type: {shape.placeholder_format.type}")

            if shape.has_text_frame and shape.text_frame.text:
                found_text = shape.text_frame.text
                found_text = found_text.replace("\r", "\n")
                print("%s" % found_text)


def is_text_placeholder(phType):
    """
    BODY = 2 # Body
    CENTER_TITLE = 3 # Center Title
    FOOTER = 15 # Footer
    HEADER = 14 # Header
    SUBTITLE = 4 # Subtitle
    TITLE = 1 # Title
    VERTICAL_BODY = 6 # Vertical Body
    VERTICAL_TITLE = 5 # Vertical Title
    """
    if phType >= PP_PLACEHOLDER_TYPE.TITLE and phType <= PP_PLACEHOLDER_TYPE.VERTICAL_BODY:
        return True

    return False


def get_placeholder_text(shapes):
    text_dict = {}
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.has_text_frame:
            phType = shape.placeholder_format.type
            if is_text_placeholder(phType):
                text = shape.text_frame.text
                text_dict[phType] = text

    return text_dict


def find_slide_by_id(slides, slide_id: int) -> tuple[int, Any]:
    for i, slide in enumerate(slides):
        if slide.slide_id == slide_id:
            return i, slide

    return -1, None


def _get_blank_slide_layout(prs):
    layout_items_count = [len(layout.placeholders) for layout in prs.slide_layouts]
    min_items = min(layout_items_count)
    blank_layout_id = layout_items_count.index(min_items)
    return prs.slide_layouts[blank_layout_id]


# https://github.com/scanny/python-pptx/issues/132
def duplicate_slide(prs, index):
    """Duplicate the slide with the given index in prs.

    Adds slide to the end of the presentation"""
    source = prs.slides[index]
    dest = prs.slides.add_slide(source.slide_layout)
    # blank_slide_layout = _get_blank_slide_layout(prs)
    # dest = prs.slides.add_slide(blank_slide_layout)
    sldIdLst = prs._element.get_or_add_sldIdLst()
    rIds = [sldId.rId for sldId in sldIdLst]
    prs.part.rename_slide_parts(rIds)

    for shape in source.shapes:
        newel = copy.deepcopy(shape.element)
        dest.shapes._spTree.insert_element_before(newel, "p:extLst")

    rels = source.part.rels._rels
    for key in rels:
        rel = rels[key]
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" not in rel.reltype:
            target = rel._target
            # if the relationship was a chart, we need to duplicate the embedded chart part and xlsx
            if "chart" in rel.reltype:
                partname = target.package.next_partname(ChartPart.partname_template)
                xlsx_blob = target.chart_workbook.xlsx_part.blob
                target = ChartPart(partname, target.content_type, copy.deepcopy(target._element), package=target.package)

                target.chart_workbook.xlsx_part = EmbeddedXlsxPart.new(xlsx_blob, target.package)

            if rel.is_external:
                dest.part.rels.get_or_add_ext_rel(rel.reltype, rel._target)
            else:
                dest.part.rels.get_or_add(rel.reltype, rel._target)

            # dest.part.rels._add_relationship(rel.reltype, target, rel.rId)

    if source.has_notes_slide:
        txt = source.notes_slide.notes_text_frame.text
        dest.notes_slide.notes_text_frame.text = txt

    if source.background and source.background._cSld.bg is not None:
        el = source.background._cSld.bg.bgPr
        newel = copy.deepcopy(el)
        cSld = dest.background._cSld
        cSld.get_or_add_bgPr()
        cSld.bg._remove_bgPr()
        cSld.bg._insert_bgPr(newel)

    return dest


# https://github.com/scanny/python-pptx/issues/68#issuecomment-575736354
def move_slide(slides, slide, new_idx):
    slides._sldIdLst.insert(new_idx, slides._sldIdLst[slides.index(slide)])


def copy_notes(source_slide, dest_slide):
    """copy_notes() fixes where duplicate_slide() may not properly copy NotesSlide placeholders.

    Note: This function is not currently used in powerpoint_pptx as duplicate_slide()
    already handles notes slide copying. It's kept for API compatibility with win32 version.
    """
    if not source_slide.has_notes_slide or not dest_slide.has_notes_slide:
        return

    # Get placeholder text from source notes slide
    src_shapes = source_slide.notes_slide.shapes
    text_dict = get_placeholder_text(src_shapes)

    # Copy placeholder text to destination notes slide
    dst_shapes = dest_slide.notes_slide.shapes
    for shape in dst_shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.has_text_frame:
            phType = shape.placeholder_format.type
            if phType in text_dict:
                shape.text_frame.text = text_dict[phType]


class Presentation(PresentationBase):
    """Presentation is a wrapper around python-ppt.P.
    All the index used in the function is 0-based and converted to 1-based while calling COM functions.
    """

    def __init__(self, app, prs):
        super().__init__(app, prs)

    def _close(self):
        # for python-pptx, no explicit close is required.
        self.reset()

    def slide_count(self):
        return len(self.prs.slides)

    def _fetch_slide_cache(self, slide_index):
        sc = SlideCache()

        slide = self.prs.slides[slide_index]
        sc.id = slide.slide_id

        sc.slide_text = get_texts_in_shapes(slide.shapes)
        sc.notes_text = get_texts_in_shapes(slide.notes_slide.shapes)

        sc.valid_cache = True

        return sc

    def slide_index_to_ID(self, var):
        if isinstance(var, int):
            slide = self.prs.slides[var]
            return slide.slide_id
        elif isinstance(var, list):
            result = []
            for index in var:
                slide = self.prs.slides[index]
                result.append(slide.slide_id)

            return result

    def slide_ID_to_index(self, var):
        if isinstance(var, int):
            index, _slide = find_slide_by_id(self.prs.slides, var)
            return index
        elif isinstance(var, list):
            result = []
            for sid in var:
                index, _slide = find_slide_by_id(self.prs.slides, sid)
                result.append(index)

            return result

    def _replace_texts_in_slide_shapes(self, slide_index, find_replace_dict):
        slide = self.prs.slides[slide_index]
        shapes = slide.shapes
        replace_texts_in_shapes(shapes, find_replace_dict)

    def duplicate_slides(self, source_location, insert_location=None, copy=1):
        source_slides_count = len(self.prs.slides)

        # validate source_location and set source_count
        # Match win32 behavior: single-item list is treated as int
        source_count = 1
        if isinstance(source_location, int):
            if not (0 <= source_location < source_slides_count):
                return 0
            source_count = 1
        elif isinstance(source_location, list):
            if len(source_location) == 1:
                # Treat single-item list as int (match win32 behavior)
                source_location = source_location[0]
                if not (0 <= source_location < source_slides_count):
                    return 0
                source_count = 1
            else:
                source_count = len(source_location)
                if source_count == 0:
                    return 0
                for src_index in source_location:
                    if not (0 <= src_index < source_slides_count):
                        return 0
        else:
            return 0

        # insert_location can change by copying slides.
        # So save the ID and use it later.
        # Match win32 behavior: when None, append at end (slide_count() - 1 in 0-based = slide_count() in 0-based for append)
        if insert_location is None:
            insert_location = self.slide_count()  # Append at end (0-based: position after last slide)
        elif isinstance(insert_location, int):
            insert_location = insert_location
        else:
            raise ValueError("Invalid insert location")

        added_count = 0

        # https://docs.microsoft.com/en-us/office/vba/api/powerpoint.sliderange.duplicate
        for _ in range(copy):
            insert_location1 = insert_location
            if isinstance(source_location, int):
                if insert_location > source_location:
                    insert_location1 = insert_location1 + source_count
            elif isinstance(source_location, list):
                # For list, check if any source location is after insert location
                max_source = max(source_location) if source_location else 0
                if insert_location > max_source:
                    insert_location1 = insert_location1 + source_count

            if isinstance(source_location, int):
                dst_slide = duplicate_slide(self.prs, source_location)
                move_slide(self.prs.slides, dst_slide, insert_location1)
            elif isinstance(source_location, list):
                # Duplicate slides in reverse order to maintain correct indices
                duplicated_slides = []
                for src_idx in reversed(source_location):
                    dst_slide = duplicate_slide(self.prs, src_idx)
                    duplicated_slides.insert(0, dst_slide)
                # Move slides to correct positions
                for i, dst_slide in enumerate(duplicated_slides):
                    move_slide(self.prs.slides, dst_slide, insert_location1 + i)

            added_count = added_count + source_count

            self._slides_inserted(insert_location - 1, source_count)

            # Update source_location for next iteration
            if isinstance(source_location, int):
                if source_location >= insert_location:
                    source_location = source_location + source_count
            elif isinstance(source_location, list):
                # Update all indices in the list
                source_location = [idx + source_count if idx >= insert_location else idx for idx in source_location]

            insert_location = insert_location + source_count

        return added_count

    @staticmethod
    def insert_template_slide(prs, index: int, slide_layout):
        slide = prs.slides.add_slide(slide_layout)
        if index != -1:
            move_slide(prs.slides, slide, index)
        return slide

    def insert_blank_slide(self, slide_index):
        layout = _get_blank_slide_layout(self.prs)
        Presentation.insert_template_slide(self.prs, slide_index, layout)

        self._slides_inserted(slide_index, 1)

    def delete_slide(self, slide_index):
        # https://github.com/scanny/python-pptx/issues/67
        # Make dictionary with necessary information
        slide = self.prs.slides[slide_index]
        id_dict = {slide.id: [i, slide.rId] for i, slide in enumerate(self.prs.slides._sldIdLst)}
        slide_id = slide.slide_id
        self.prs.part.drop_rel(id_dict[slide_id][1])
        del self.prs.slides._sldIdLst[id_dict[slide_id][0]]

        self._slides_deleted(slide_index, 1)

    def copy_all_and_close(self):
        """Store the presentation in the app for later pasting, then close this presentation."""
        # Store the source presentation in the app for later pasting
        self.app._clipboard_presentation = self.prs
        self.close()

    def _paste_keep_source_formatting(self, insert_location):
        """Paste slides from the clipboard presentation (stored in copy_all_and_close) into this presentation."""
        if not hasattr(self.app, "_clipboard_presentation") or self.app._clipboard_presentation is None:
            return

        source_prs = self.app._clipboard_presentation
        source_slide_count = len(source_prs.slides)

        if source_slide_count == 0:
            self.app._clipboard_presentation = None
            return

        # Copy all slides from source to destination
        # We need to insert them in reverse order to maintain correct positions
        inserted_slides = []
        for i in range(source_slide_count):
            source_slide = source_prs.slides[i]
            # Get the layout from source slide or use blank layout
            try:
                layout = source_slide.slide_layout
            except:
                layout = _get_blank_slide_layout(self.prs)

            # Create new slide (will be added at the end initially)
            dest_slide = self.prs.slides.add_slide(layout)
            inserted_slides.append((dest_slide, source_slide))

        # Now move slides to correct positions and copy content
        for idx, (dest_slide, source_slide) in enumerate(inserted_slides):
            target_idx = insert_location + idx
            # Move to correct position
            if target_idx < len(self.prs.slides) - 1:
                move_slide(self.prs.slides, dest_slide, target_idx)

            # Copy all shapes from source to destination
            for shape in source_slide.shapes:
                newel = copy.deepcopy(shape.element)
                dest_slide.shapes._spTree.insert_element_before(newel, "p:extLst")

            # Copy relationships
            rels = source_slide.part.rels._rels
            for key in rels:
                rel = rels[key]
                if "notesSlide" not in rel.reltype:
                    target = rel._target
                    if "chart" in rel.reltype:
                        partname = target.package.next_partname(ChartPart.partname_template)
                        xlsx_blob = target.chart_workbook.xlsx_part.blob
                        target = ChartPart(partname, target.content_type, copy.deepcopy(target._element), package=target.package)
                        target.chart_workbook.xlsx_part = EmbeddedXlsxPart.new(xlsx_blob, target.package)

                    if rel.is_external:
                        dest_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel._target)
                    else:
                        dest_slide.part.rels.get_or_add(rel.reltype, rel._target)

            # Copy notes slide
            if source_slide.has_notes_slide and dest_slide.has_notes_slide:
                txt = source_slide.notes_slide.notes_text_frame.text
                dest_slide.notes_slide.notes_text_frame.text = txt

            # Copy background
            if source_slide.background and source_slide.background._cSld.bg is not None:
                el = source_slide.background._cSld.bg.bgPr
                newel = copy.deepcopy(el)
                cSld = dest_slide.background._cSld
                cSld.get_or_add_bgPr()
                cSld.bg._remove_bgPr()
                cSld.bg._insert_bgPr(newel)

        # Clean up clipboard
        self.app._clipboard_presentation = None

    def saveas(self, filename):
        """Save .pptx as files using python-pptx."""
        self.prs.save(filename)

    def saveas_format(self, filename, image_type="png"):
        """Save .pptx as image files.

        Note: python-pptx does not support exporting presentations as image formats.
        This method raises NotImplementedError. Use export_slide_as for individual slides.
        """
        raise NotImplementedError(
            "python-pptx does not support exporting presentations as image formats. Use export_slide_as() to export individual slides."
        )

    def export_slide_as(self, slideno, filename, image_type="png"):
        """Save all shapes in slide as image files.

        Note: python-pptx does not support rendering slides as images.
        This method raises NotImplementedError. python-pptx is a library for creating/editing
        PPTX files, not for rendering them. To export slides as images, you would need:
        - PowerPoint application (COM interface on Windows, AppleScript on macOS)
        - Or a rendering library like LibreOffice in headless mode
        """
        raise NotImplementedError(
            "python-pptx does not support exporting slides as images. "
            "python-pptx is a library for creating/editing PPTX files, not rendering them. "
            "To export slides as images, use powerpoint_win32 or powerpoint_osx modules."
        )

    def export_shapes_as(self, slideno, filename, image_type="png"):
        """Save all shapes in slide as image file.
        It only exports the shapes not slide itself and the dimension may not what you expect,
        unless the whole selected shapes cover the slide.
        This method exports transparent images in good condition compared to the above saveas_format.

        Note: python-pptx does not support rendering shapes as images.
        This method raises NotImplementedError. python-pptx is a library for creating/editing
        PPTX files, not for rendering them.

        https://stackoverflow.com/questions/5713676/ppt-to-png-with-transparent-background
        https://github.com/DefinitelyTyped/DefinitelyTyped/blob/a102789c764788888edc6a89542cd90f08fdce3d/types/activex-powerpoint/index.d.ts#L4234
        https://stackoverflow.com/questions/45299899/graphic-quality-in-powerpoint-graphics-export
        """
        raise NotImplementedError(
            "python-pptx does not support exporting shapes as images. "
            "python-pptx is a library for creating/editing PPTX files, not rendering them. "
            "To export shapes as images, use powerpoint_win32 or powerpoint_osx modules."
        )


class App(PPTAppBase):
    @staticmethod
    def is_running():
        return True

    def __init__(self):
        self.presentation = PptxPres()
        self._clipboard_presentation = None  # Used for copy/paste operations

    def new_presentation(self):
        layout = _get_blank_slide_layout(self.presentation)
        _slide = Presentation.insert_template_slide(self.presentation, 0, layout)

        return Presentation(self, self.presentation)

    def open_presentation(self, filename):
        with open(filename, "rb") as f:
            source_stream = BytesIO(f.read())
            prs = PptxPres(source_stream)
            source_stream.close()

        if prs is None:
            return None

        return Presentation(self, prs)

    def quit(self, force=False, only_if_empty=True):
        """Quit the application.

        Note: For python-pptx, there is no actual application to quit,
        so this method does nothing. It's provided for API compatibility with win32/osx.
        """
        # for python-pptx, there is no app, so no quitting.
        pass

    def quit_if_empty(self):
        """Quit if no presentations are open.

        Note: For python-pptx, there is no actual application to quit,
        so this method does nothing. It's provided for API compatibility.
        """
        # for python-pptx, there is no app, so no quitting.
        pass
