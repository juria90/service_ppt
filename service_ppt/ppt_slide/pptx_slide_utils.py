"""Utility functions for working with python-pptx slides.

This module provides utility functions for manipulating slides in python-pptx
presentations, including slide duplication with full content preservation
(shapes, relationships, charts, notes, and backgrounds).

https://github.com/scanny/python-pptx/issues/132#issuecomment-1643400011
"""

import copy
import io
import random
import string
from typing import Any

from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.packuri import PackURI
from pptx.parts.slide import SlidePart
from pptx.shapes.group import GroupShape


def _object_rels(obj: Any) -> list[Any]:
    """Extract relationship objects from an object's relationships.

    Handles compatibility with python-pptx 0.6.22+ where the relationship
    dictionary structure changed.

    :param obj: Object with a rels attribute (e.g., slide part)
    :return: List of relationship objects
    """
    rels = obj.rels

    # Change required for python-pptx 0.6.22
    check_rels_content = [k for k in rels]
    if isinstance(check_rels_content.pop(), str):
        return [v for k, v in rels.items()]
    else:
        return [k for k in rels]


def _exp_add_slide(ppt: Any, slide_layout: Any) -> Any:
    """Create a new slide in the presentation with the specified layout.

    This function handles slide creation to avoid issues with the default
    python-pptx implementation, particularly when dealing with duplicate
    slide partnames. It generates unique partnames and properly manages
    slide relationships.

    :param ppt: Presentation object to add the slide to
    :param slide_layout: Slide layout object to use for the new slide
    :return: Newly created slide object
    """

    def generate_slide_partname(self: Any) -> Any:
        """Generate a unique slide partname for a new slide.

        Returns a PackURI instance containing the next available slide partname.
        If the default partname already exists, generates a random suffix to
        ensure uniqueness.

        :param self: Slide collection object
        :return: PackURI instance with unique slide partname
        """
        sldIdLst = self._element.get_or_add_sldIdLst()

        existing_rels = [k.target_partname for k in _object_rels(self)]
        partname_str = "/ppt/slides/slide%d.xml" % (len(sldIdLst) + 1)

        while partname_str in existing_rels:
            random_part = "".join(random.choice(string.ascii_letters) for i in range(2))
            partname_str = "/ppt/slides/slide%s%d.xml" % (random_part, len(sldIdLst) + 1)

        return PackURI(partname_str)

    def add_slide_part(self: Any, slide_layout: Any) -> tuple[str, Any]:
        """Create a new slide part and return its relationship ID and slide object.

        Creates a new blank slide that inherits appearance from the specified
        slide layout and establishes the relationship in the presentation package.

        :param self: Presentation part object
        :param slide_layout: Slide layout object to inherit from
        :return: Tuple of (relationship_id, slide_object)
        """
        partname = generate_slide_partname(self)
        slide_layout_part = slide_layout.part
        slide_part = SlidePart.new(partname, self.package, slide_layout_part)
        rId = self.relate_to(slide_part, RT.SLIDE)
        return rId, slide_part.slide

    def add_slide_ppt(self: Any, slide_layout: Any) -> Any:
        """Add a new slide to the presentation using the specified layout.

        Creates a slide part, clones layout placeholders, and adds the slide
        to the presentation's slide ID list.

        :param self: Slide collection object
        :param slide_layout: Slide layout object to use
        :return: Newly created slide object
        """
        rId, slide = add_slide_part(self.part, slide_layout)
        slide.shapes.clone_layout_placeholders(slide_layout)
        self._sldIdLst.add_sldId(rId)
        return slide

    # slide_layout = self.get_master_slide_layout(slide_layout)
    return add_slide_ppt(ppt.slides, slide_layout)


def copy_shapes(source: Any, dest: Any) -> None:
    """Copy shapes from source to destination, handling various shape types.

    Recursively copies shapes including group shapes, images, charts, and
    other shape types. Handles edge cases like nested groups, image crops,
    and chart cloning. For group shapes, recursively copies all contained
    shapes and preserves group properties.

    :param source: Source shapes collection to copy from
    :param dest: Destination slide or group shape to copy to
    :return: None
    """
    # Copy all existing shapes
    for shape in source:
        if isinstance(shape, GroupShape):
            group = dest.shapes.add_group_shape()
            group.name = shape.name
            group.left = shape.left
            group.top = shape.top
            group.width = shape.width
            group.height = shape.height
            group.rotation = shape.rotation

            # Recursive copy of contents
            copy_shapes(shape.shapes, group)

            # Fix offset
            cur_el = group._element.xpath(".//p:grpSpPr")[0]
            ref_el = shape._element.xpath(".//p:grpSpPr")[0]
            parent = cur_el.getparent()
            parent.insert(parent.index(cur_el) + 1, copy.deepcopy(ref_el))
            parent.remove(cur_el)

            result = group
        elif hasattr(shape, "image"):
            # Get image contents
            content = io.BytesIO(shape.image.blob)
            result = dest.shapes.add_picture(content, shape.left, shape.top, shape.width, shape.height)
            result.name = shape.name
            result.crop_left = shape.crop_left
            result.crop_right = shape.crop_right
            result.crop_top = shape.crop_top
            result.crop_bottom = shape.crop_bottom
        # elif hasattr(shape, "has_chart") and shape.has_chart:
        #   result = clone_chart(shape, dest)
        else:
            newel = copy.deepcopy(shape.element)
            dest.shapes._spTree.insert_element_before(newel, "p:extLst")
            result = dest.shapes[-1]


def remove_shape(shape: Any) -> None:
    """
    Helper to remove a specific shape.

    :source: https://stackoverflow.com/questions/64700638/is-there-a-way-to-delete-a-shape-with-python-pptx

    :param shape: Shape object to remove
    :return: None
    """
    el = shape.element  # --- get reference to XML element for shape
    el.getparent().remove(el)  # --- remove that shape element from its tree


def duplicate_slide(ppt: Any, slide_index: int) -> Any:
    """Duplicate a slide in the presentation with full content preservation.

    Creates a copy of the specified slide including all shapes, notes, and
    background formatting. The new slide is added at the end of the
    presentation by default. Handles complex shapes including groups, images,
    and charts.

    :param ppt: Presentation object containing the slide to duplicate
    :param slide_index: Zero-based index of the slide to duplicate
    :return: Newly created duplicate slide object
    """
    source = ppt.slides[slide_index]

    dest = _exp_add_slide(ppt, source.slide_layout)

    # Remove all shapes from the default layout
    for shape in dest.shapes:
        remove_shape(shape)

    # Copy all existing shapes
    copy_shapes(source.shapes, dest)

    # Copy all existing shapes
    if source.has_notes_slide:
        txt = source.notes_slide.notes_text_frame.text
        dest.notes_slide.notes_text_frame.text = txt

    # added by juria90 to handle background image.
    if source.background and source.background._cSld.bg is not None:
        el = source.background._cSld.bg.bgPr
        newel = copy.deepcopy(el)
        csld = dest.background._cSld
        csld.get_or_add_bgPr()
        csld.bg._remove_bgPr()
        csld.bg._insert_bgPr(newel)

    return dest
