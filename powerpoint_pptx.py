"""powerpoint_win32.py implements App and Presentation for Powerpoint in Windows using COM interface.
"""

import copy
from enum import IntEnum
from io import BytesIO
import re
from typing import Any, Callable, Tuple

from pptx import Presentation as PptxPres
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.parts.chart import ChartPart
from pptx.parts.embeddedpackage import EmbeddedXlsxPart

from powerpoint_base import SlideCache, PresentationBase, PPTAppBase


class SlideLayout(IntEnum):
    SLD_LAYOUT_TITLE = 0
    SLD_LAYOUT_TITLE_AND_CONTENT = 1
    SLD_LAYOUT_SECTION_HEADER = 2
    SLD_LAYOUT_TWO_CONTENT = 3
    SLD_LAYOUT_COMPARISON = 4
    SLD_LAYOUT_TITLE_ONLY = 5
    SLD_LAYOUT_BLANK = 6
    SLD_LAYOUT_CONTENT_WITH_CAPTION = 7
    SLD_LAYOUT_PICTURE_WITH_CAPTION = 8


def iterate_text_objects_in_shapes(shapes, fn: Callable = None):
    for shape in shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if fn:
                        try:
                            result = fn(run)
                            if result:
                                yield run
                        except StopIteration:
                            return
                    else:
                        yield run

        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if fn:
                        try:
                            result = fn(run)
                            if result:
                                yield run
                        except StopIteration:
                            return
                    else:
                        yield run

        if shape.has_chart:
            chart = shape.chart
            if chart.has_title:
                chart_title = chart.chart_title
                if chart_title.has_text_frame:
                    text_frame = chart_title.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            if fn:
                                try:
                                    result = fn(run)
                                    if result:
                                        yield run
                                except StopIteration:
                                    return
                            else:
                                yield run


def get_texts_in_shapes(shapes):
    texts = [obj.text for obj in iterate_text_objects_in_shapes(shapes)]
    return texts


def find_text_in_shapes(shapes, text, match_case=True, whole_words=True):
    # https://stackoverflow.com/questions/25483114/regex-to-find-whole-word-in-text-but-case-insensitive
    pattern = ""
    if match_case is False:
        pattern += r"(?i)"
    if whole_words:
        pattern += r"\b" + re.escape(text) + r"\b"
    else:
        pattern += re.escape(text)

    re_text = re.compile(pattern)

    def match_string(obj):
        nonlocal re_text
        return re_text.search(obj.text) is not None

    texts = [obj.text for obj in iterate_text_objects_in_shapes(shapes, match_string)]
    for shape in shapes:
        if shape.has_text_frame:
            found_text = text in shape.text
            if found_text is not None:
                return found_text

    return texts[0] if len(texts) > 0 else None


def replace_texts_in_shapes(shapes, text_dict):
    count = 0

    def replace_texts_in_run(obj):
        nonlocal count, text_dict
        for from_text, to_text in text_dict.items():
            if (obj.text.find(from_text)) != -1:
                obj.text = obj.text.replace(from_text, to_text)
                count = count + 1

    for _ in iterate_text_objects_in_shapes(shapes, replace_texts_in_run):
        pass

    return count


def dump_shapes(shapes):
    for i, shape in enumerate(shapes):
        if shape.has_text_frame:

            print("shape %d: %d" % (i, shape.Type))

            if shape.Type == MSO_SHAPE_TYPE.PLACEHOLDER:
                print(f"PlaceholderFormat.Type: {shape.placeholder_format.type}")

            if shape.text:
                found_text = shape.text
                found_text = found_text.replace("\r", "\n")
                print("%s" % found_text)


def is_text_placeholder(phType):
    """
    BODY = 2 # Body
    CENTER_TITLE = 3 # Center Title
    FOOTER = 15 # Footer
    HEADER = 14 # Header
    SUBTITLE = 4 # Subtitle
    TITLE = 1 # Title
    VERTICAL_BODY = 6 # Vertical Body
    VERTICAL_TITLE = 5 # Vertical Title
    """
    if PP_PLACEHOLDER.TITLE <= phType and phType <= PP_PLACEHOLDER.VERTICAL_BODY:
        return True

    return False


def get_placeholder_text(shapes):
    text_dict = {}
    for shape in shapes:

        if shape.type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.text:
            phType = shape.placeholder_format.type
            if is_text_placeholder(phType):
                text_range = shape.TextFrame.TextRange
                text = text_range.Text
                text_dict[phType] = text

    return text_dict


def find_slide_by_id(slides, slide_id: int) -> Tuple[int, Any]:
    for i, slide in enumerate(slides):
        if slide.slide_id == slide_id:
            return i, slide

    return -1, None


def _get_blank_slide_layout(prs):
    layout_items_count = [len(layout.placeholders) for layout in prs.slide_layouts]
    min_items = min(layout_items_count)
    blank_layout_id = layout_items_count.index(min_items)
    return prs.slide_layouts[blank_layout_id]


# https://github.com/scanny/python-pptx/issues/132
def duplicate_slide(prs, index):
    """Duplicate the slide with the given index in prs.

    Adds slide to the end of the presentation"""
    source = prs.slides[index]
    dest = prs.slides.add_slide(source.slide_layout)
    # blank_slide_layout = _get_blank_slide_layout(prs)
    # dest = prs.slides.add_slide(blank_slide_layout)
    sldIdLst = prs._element.get_or_add_sldIdLst()
    rIds = [sldId.rId for sldId in sldIdLst]
    prs.part.rename_slide_parts(rIds)

    for shape in source.shapes:
        newel = copy.deepcopy(shape.element)
        dest.shapes._spTree.insert_element_before(newel, "p:extLst")

    rels = source.part.rels._rels
    for key in rels:
        rel = rels[key]
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" not in rel.reltype:
            target = rel._target
            # if the relationship was a chart, we need to duplicate the embedded chart part and xlsx
            if "chart" in rel.reltype:
                partname = target.package.next_partname(ChartPart.partname_template)
                xlsx_blob = target.chart_workbook.xlsx_part.blob
                target = ChartPart(partname, target.content_type, copy.deepcopy(target._element), package=target.package)

                target.chart_workbook.xlsx_part = EmbeddedXlsxPart.new(xlsx_blob, target.package)

            if rel.is_external:
                dest.part.rels.get_or_add_ext_rel(rel.reltype, rel._target)
            else:
                dest.part.rels.get_or_add(rel.reltype, rel._target)

            # dest.part.rels._add_relationship(rel.reltype, target, rel.rId)

    if source.has_notes_slide:
        txt = source.notes_slide.notes_text_frame.text
        dest.notes_slide.notes_text_frame.text = txt

    if source.background:
        el = source.background._cSld.bg.bgPr
        newel = copy.deepcopy(el)
        cSld = dest.background._cSld
        cSld.get_or_add_bgPr()
        cSld.bg._remove_bgPr()
        cSld.bg._insert_bgPr(newel)

    return dest


# https://github.com/scanny/python-pptx/issues/68#issuecomment-575736354
def move_slide(slides, slide, new_idx):
    slides._sldIdLst.insert(new_idx, slides._sldIdLst[slides.index(slide)])


def copy_notes(source_range, dst_range):
    """copy_notes() fixes where SlideRange.Duplicate() does not copy NotesSlide."""
    # print('source_range')
    src_shapes = source_range.NotesPage.Shapes
    text_dict = get_placeholder_text(src_shapes)
    # dump_shapes(src_shapes)

    # print('dst_range')
    shapes = dst_range.NotesPage.Shapes
    # dump_shapes(shapes)
    for i in range(shapes.Count):
        shape = shapes.Item(i + 1)

        if shape.type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.text:
            phType = shape.placeholder_format.type
            if phType in text_dict:
                shape.TextFrame.TextRange.Text = text_dict[phType]

    # print('dst_range after change')
    # shapes = dst_range.NotesPage.Shapes
    # dump_shapes(shapes)


class Presentation(PresentationBase):
    """Presentation is a wrapper around python-ppt.P.
    All the index used in the function is 0-based and converted to 1-based while calling COM functions.
    """

    def __init__(self, app, prs):
        super().__init__(app, prs)

    def slide_count(self):
        return len(self.prs.slides)

    def _fetch_slide_cache(self, slide_index):
        sc = SlideCache()

        slide = self.prs.slides[slide_index]
        sc.id = slide.slide_id

        sc.slide_text = get_texts_in_shapes(slide.shapes)
        sc.notes_text = get_texts_in_shapes(slide.notes_slide.shapes)

        sc.valid_cache = True

        return sc

    def slide_index_to_ID(self, var):
        if isinstance(var, int):
            slide = self.prs.slides[var]
            return slide.slide_id
        elif isinstance(var, list):
            result = []
            for index in var:
                slide = self.prs.slides[index]
                result.append(slide.slide_id)

            return result

    def slide_ID_to_index(self, var):
        if isinstance(var, int):
            index, _slide = find_slide_by_id(self.prs.slides, var)
            return index
        elif isinstance(var, list):
            result = []
            for sid in var:
                index, _slide = find_slide_by_id(self.prs.slides, sid)
                result.append(index)

            return result

    def _replace_texts_in_slide_shapes(self, slide_index, find_replace_dict):
        slide = self.prs.slides[slide_index]
        shapes = slide.shapes
        replace_texts_in_shapes(shapes, find_replace_dict)

    def duplicate_slides(self, source_location, insert_location=None, copy=1):

        source_slides_count = len(self.prs.slides)

        # validate source_location and set source_count
        source_count = 1
        if isinstance(source_location, int):
            if not (0 <= source_location or source_location < source_slides_count):
                return 0

            source_count = 1
        elif isinstance(source_location, list):
            source_count = len(source_location)
            if source_count == 0:
                return 0

            for src_index in source_location:
                if not (0 <= src_index or src_index < source_slides_count):
                    return 0
        else:
            return 0

        # insert_location can change by copying slides.
        # So save the ID and use it later.
        if insert_location is None:
            insert_location = self.slide_count()
        elif isinstance(insert_location, int):
            insert_location = insert_location
        else:
            raise ValueError("Invalid insert location")

        added_count = 0

        # https://docs.microsoft.com/en-us/office/vba/api/powerpoint.sliderange.duplicate
        for _ in range(copy):
            insert_location1 = insert_location
            if insert_location > source_location:
                insert_location1 = insert_location1 + source_count

            if isinstance(source_location, int):
                dst_slide = duplicate_slide(self.prs, source_location)
                move_slide(self.prs.slides, dst_slide, insert_location1)
            elif isinstance(source_location, list):
                assert False
                pass

            added_count = added_count + source_count

            self._slides_inserted(insert_location - 1, source_count)
            if source_location >= insert_location:
                source_location = source_location + source_count

            insert_location = insert_location + source_count

        return added_count

    @staticmethod
    def insert_template_slide(prs, index: int, slide_layout):
        slide = prs.slides.add_slide(slide_layout)
        if index != -1:
            move_slide(prs.slides, slide, index)
        return slide

    def insert_blank_slide(self, slide_index):
        layout = _get_blank_slide_layout(self.prs)
        Presentation.insert_template_slide(self.prs, slide_index, layout)

        self._slides_inserted(slide_index, 1)

    def delete_slide(self, slide_index):
        # https://github.com/scanny/python-pptx/issues/67
        # Make dictionary with necessary information
        slide = self.prs.slides[slide_index]
        id_dict = {slide.id: [i, slide.rId] for i, slide in enumerate(self.prs.slides._sldIdLst)}
        slide_id = slide.slide_id
        self.prs.part.drop_rel(id_dict[slide_id][1])
        del self.prs.slides._sldIdLst[id_dict[slide_id][0]]

        self._slides_deleted(slide_index, 1)

    def copy_all_and_close(self):
        return

        all_range = self.prs.slides.Range()
        all_range.Copy()
        self.close()

    def _paste_keep_source_formatting(self, insert_location):
        pass

    def saveas(self, filename):
        """Save .pptx as files using Powerpoint.Application COM service."""
        self.prs.save(filename)

    def saveas_format(self, filename, image_type="png"):
        """Save .pptx as files using Powerpoint.Application COM service."""

    def export_slide_as(self, slideno, filename, image_type="png"):
        """Save all shapes in slide as image files."""

    def export_shapes_as(self, slideno, filename, image_type="png"):
        """Save all shapes in slide as image file.
        It only exports the shapes not slide itself and the dimension may not what you expect,
        unless the whole selected shapes cover the slide.
        This method exports transparent images in good condition compared to the above saveas_format.

        https://stackoverflow.com/questions/5713676/ppt-to-png-with-transparent-background
        https://github.com/DefinitelyTyped/DefinitelyTyped/blob/a102789c764788888edc6a89542cd90f08fdce3d/types/activex-powerpoint/index.d.ts#L4234
        https://stackoverflow.com/questions/45299899/graphic-quality-in-powerpoint-graphics-export
        """

    def close(self):
        # for python-pptx, no explicit close is required.
        self.reset()


class App(PPTAppBase):
    @staticmethod
    def is_running():
        return True

    def __init__(self):
        self.presentation = PptxPres()

    def new_presentation(self):
        layout = _get_blank_slide_layout(self.presentation)
        _slide = Presentation.insert_template_slide(self.presentation, 0, layout)

        return Presentation(self, self.presentation)

    def open_presentation(self, filename):
        with open(filename, "rb") as f:
            source_stream = BytesIO(f.read())
            prs = PptxPres(source_stream)
            source_stream.close()

        if prs is None:
            return None

        return Presentation(self, prs)

    def quit_if_empty(self):
        # for python-pptx, there is no app, so no quitting.
        pass
