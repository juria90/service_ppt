"""Tests for powerpoint_pptx module.

This module contains unit tests for the Presentation class in powerpoint_pptx,
testing major features like slide manipulation, text replacement, and file operations.
"""

from pathlib import Path
import tempfile

import pytest
from pptx import Presentation as PptxPres

from service_ppt.ppt_slide.powerpoint_pptx import App, Presentation


@pytest.fixture
def app():
    """Create an App instance for testing."""
    return App()


@pytest.fixture
def empty_presentation(app):
    """Create an empty presentation for testing."""
    return app.new_presentation()


@pytest.fixture
def presentation_with_slides(app):
    """Create a presentation with multiple slides for testing."""
    prs = app.new_presentation()

    # Delete the initial blank slide and create new ones
    if prs.slide_count() > 0:
        prs.delete_slide(0)

    # Add a few slides with text
    for i in range(3):
        prs.insert_blank_slide(i)
        slide = prs.prs.slides[i]
        # Add a text box with some text
        textbox = slide.shapes.add_textbox(0, 0, 100, 100)
        textbox.text_frame.text = f"Slide {i} content"

    return prs


class TestPresentationBasic:
    """Test basic Presentation functionality."""

    def test_slide_count_empty(self, empty_presentation):
        """Test slide_count() on empty presentation."""
        assert empty_presentation.slide_count() == 1  # New presentation has 1 blank slide

    def test_slide_count_with_slides(self, presentation_with_slides):
        """Test slide_count() with multiple slides."""
        assert presentation_with_slides.slide_count() == 3

    def test_slide_index_to_ID_single(self, presentation_with_slides):
        """Test slide_index_to_ID() with single index."""
        slide_id = presentation_with_slides.slide_index_to_ID(0)
        assert isinstance(slide_id, int)
        assert slide_id > 0

    def test_slide_index_to_ID_list(self, presentation_with_slides):
        """Test slide_index_to_ID() with list of indices."""
        slide_ids = presentation_with_slides.slide_index_to_ID([0, 1, 2])
        assert isinstance(slide_ids, list)
        assert len(slide_ids) == 3
        assert all(isinstance(sid, int) and sid > 0 for sid in slide_ids)
        # IDs should be unique
        assert len(set(slide_ids)) == 3

    def test_slide_ID_to_index_single(self, presentation_with_slides):
        """Test slide_ID_to_index() with single ID."""
        slide_id = presentation_with_slides.slide_index_to_ID(0)
        index = presentation_with_slides.slide_ID_to_index(slide_id)
        assert index == 0

    def test_slide_ID_to_index_list(self, presentation_with_slides):
        """Test slide_ID_to_index() with list of IDs."""
        slide_ids = presentation_with_slides.slide_index_to_ID([0, 1, 2])
        indices = presentation_with_slides.slide_ID_to_index(slide_ids)
        assert isinstance(indices, list)
        assert len(indices) == 3
        assert indices == [0, 1, 2]

    def test_fetch_slide_cache(self, presentation_with_slides):
        """Test _fetch_slide_cache() retrieves slide information."""
        cache = presentation_with_slides._fetch_slide_cache(0)
        assert cache.id > 0
        assert cache.valid_cache is True
        assert isinstance(cache.slide_text, list)
        assert isinstance(cache.notes_text, list)
        # Should have text from the slide
        assert len(cache.slide_text) > 0


class TestPresentationTextOperations:
    """Test text replacement and retrieval operations."""

    def test_replace_texts_in_slide_shapes(self, presentation_with_slides):
        """Test _replace_texts_in_slide_shapes() replaces text correctly."""
        find_replace = {"Slide 0": "Replaced Text"}
        presentation_with_slides._replace_texts_in_slide_shapes(0, find_replace)

        # Verify text was replaced
        cache = presentation_with_slides._fetch_slide_cache(0)
        assert "Replaced Text" in " ".join(cache.slide_text)

    def test_replace_texts_multiple_replacements(self, presentation_with_slides):
        """Test _replace_texts_in_slide_shapes() with multiple replacements."""
        find_replace = {"Slide 0": "First Replacement", "content": "Second Replacement"}
        presentation_with_slides._replace_texts_in_slide_shapes(0, find_replace)

        cache = presentation_with_slides._fetch_slide_cache(0)
        text = " ".join(cache.slide_text)
        assert "First Replacement" in text or "Second Replacement" in text


class TestPresentationSlideManipulation:
    """Test slide manipulation operations."""

    def test_insert_blank_slide(self, empty_presentation):
        """Test insert_blank_slide() adds a new blank slide."""
        initial_count = empty_presentation.slide_count()
        empty_presentation.insert_blank_slide(1)
        assert empty_presentation.slide_count() == initial_count + 1

    def test_insert_blank_slide_at_beginning(self, presentation_with_slides):
        """Test insert_blank_slide() at the beginning."""
        initial_count = presentation_with_slides.slide_count()
        presentation_with_slides.insert_blank_slide(0)
        assert presentation_with_slides.slide_count() == initial_count + 1

    def test_delete_slide(self, presentation_with_slides):
        """Test delete_slide() removes a slide."""
        initial_count = presentation_with_slides.slide_count()
        presentation_with_slides.delete_slide(0)
        assert presentation_with_slides.slide_count() == initial_count - 1

    def test_delete_slide_middle(self, presentation_with_slides):
        """Test delete_slide() removes slide from middle."""
        initial_count = presentation_with_slides.slide_count()
        # Get ID of middle slide before deletion
        middle_id = presentation_with_slides.slide_index_to_ID(1)
        presentation_with_slides.delete_slide(1)
        assert presentation_with_slides.slide_count() == initial_count - 1
        # Verify the slide with that ID is gone (returns -1 for not found)
        result = presentation_with_slides.slide_ID_to_index(middle_id)
        assert result == -1

    def test_duplicate_slides_single_slide(self, presentation_with_slides):
        """Test duplicate_slides() duplicates a single slide."""
        initial_count = presentation_with_slides.slide_count()
        added = presentation_with_slides.duplicate_slides(0, insert_location=1)
        assert added == 1
        assert presentation_with_slides.slide_count() == initial_count + 1

    def test_duplicate_slides_multiple_copies(self, presentation_with_slides):
        """Test duplicate_slides() creates multiple copies."""
        initial_count = presentation_with_slides.slide_count()
        added = presentation_with_slides.duplicate_slides(0, insert_location=1, copy=3)
        assert added == 3
        assert presentation_with_slides.slide_count() == initial_count + 3

    def test_duplicate_slides_list_of_slides(self, presentation_with_slides):
        """Test duplicate_slides() with list of source slides."""
        initial_count = presentation_with_slides.slide_count()
        added = presentation_with_slides.duplicate_slides([0, 2], insert_location=1)
        assert added == 2
        assert presentation_with_slides.slide_count() == initial_count + 2

    def test_duplicate_slides_append_at_end(self, presentation_with_slides):
        """Test duplicate_slides() appends at end when insert_location is None."""
        initial_count = presentation_with_slides.slide_count()
        added = presentation_with_slides.duplicate_slides(0, insert_location=None)
        assert added == 1
        assert presentation_with_slides.slide_count() == initial_count + 1

    def test_duplicate_slides_invalid_source(self, presentation_with_slides):
        """Test duplicate_slides() with invalid source location."""
        initial_count = presentation_with_slides.slide_count()
        added = presentation_with_slides.duplicate_slides(999, insert_location=0)
        assert added == 0
        assert presentation_with_slides.slide_count() == initial_count

    def test_duplicate_slides_preserves_content(self, presentation_with_slides):
        """Test duplicate_slides() preserves slide content."""
        # Get original slide text from slide 0
        original_cache = presentation_with_slides._fetch_slide_cache(0)
        original_text = " ".join(original_cache.slide_text)

        # Store the original slide count
        original_count = presentation_with_slides.slide_count()

        # Duplicate slide 0 and insert at position 1 (after slide 0)
        # Note: Due to the current implementation logic, when duplicating a slide
        # that comes before the insert location, the insert location is adjusted.
        # So inserting at position 1 actually results in the duplicate at position 2.
        presentation_with_slides.duplicate_slides(0, insert_location=1)

        # After duplication, we should have one more slide
        assert presentation_with_slides.slide_count() == original_count + 1

        # Find the duplicated slide (it should have the same content as slide 0)
        # The duplicate will be at position 2 due to the insert_location adjustment logic
        duplicated_cache = presentation_with_slides._fetch_slide_cache(2)
        duplicated_text = " ".join(duplicated_cache.slide_text)
        assert original_text == duplicated_text, f"Expected '{original_text}' but got '{duplicated_text}'"


class TestPresentationCopyPaste:
    """Test copy and paste operations."""

    def test_copy_all_and_close(self, presentation_with_slides):
        """Test copy_all_and_close() stores presentation in clipboard."""
        app = presentation_with_slides.app
        assert app._clipboard_presentation is None
        clipboard_prs = presentation_with_slides.prs
        presentation_with_slides.copy_all_and_close()
        # Check clipboard before app is reset
        assert app._clipboard_presentation is not None
        assert app._clipboard_presentation == clipboard_prs
        # After close, prs should be None (reset is called)
        assert presentation_with_slides.prs is None

    def test_paste_keep_source_formatting(self, app):
        """Test _paste_keep_source_formatting() pastes slides from clipboard."""
        # Create source presentation
        source_prs = app.new_presentation()
        # Delete the initial blank slide to start fresh
        if source_prs.slide_count() > 0:
            source_prs.delete_slide(0)
        # Add 2 slides
        source_prs.insert_blank_slide(0)
        source_prs.insert_blank_slide(1)
        source_slide_count = source_prs.slide_count()  # Should be 2

        # Copy source
        source_prs.copy_all_and_close()

        # Create destination presentation
        dest_prs = app.new_presentation()
        # Delete the initial blank slide to start fresh
        if dest_prs.slide_count() > 0:
            dest_prs.delete_slide(0)
        # Add 1 slide
        dest_prs.insert_blank_slide(0)
        initial_count = dest_prs.slide_count()  # Should be 1

        # Paste at position 1 (after the first slide)
        dest_prs._paste_keep_source_formatting(1)
        # Verify slides were pasted
        # Note: _paste_keep_source_formatting adds slides at the end first, then moves them
        # So the actual count should be initial + source
        final_count = dest_prs.slide_count()
        assert final_count >= initial_count + source_slide_count, (
            f"Expected at least {initial_count + source_slide_count} slides but got {final_count}"
        )
        assert dest_prs.app._clipboard_presentation is None  # Clipboard should be cleared

    def test_paste_keep_source_formatting_empty_clipboard(self, empty_presentation):
        """Test _paste_keep_source_formatting() with empty clipboard."""
        initial_count = empty_presentation.slide_count()
        empty_presentation._paste_keep_source_formatting(0)
        # Should not change slide count
        assert empty_presentation.slide_count() == initial_count


class TestPresentationFileOperations:
    """Test file save operations."""

    def test_saveas(self, presentation_with_slides):
        """Test saveas() saves presentation to file."""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            presentation_with_slides.saveas(tmp_path)
            # Verify file was created
            assert Path(tmp_path).exists()
            # Verify it's a valid PPTX file by opening it
            loaded_prs = PptxPres(tmp_path)
            assert loaded_prs is not None
            assert len(loaded_prs.slides) == presentation_with_slides.slide_count()
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    def test_saveas_preserves_content(self, presentation_with_slides):
        """Test saveas() preserves slide content."""
        # Add some text to a slide
        presentation_with_slides._replace_texts_in_slide_shapes(0, {"Slide 0": "Test Content"})

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            presentation_with_slides.saveas(tmp_path)
            # Load and verify content
            loaded_prs = PptxPres(tmp_path)
            loaded_presentation = Presentation(presentation_with_slides.app, loaded_prs)
            cache = loaded_presentation._fetch_slide_cache(0)
            assert "Test Content" in " ".join(cache.slide_text)
        finally:
            Path(tmp_path).unlink(missing_ok=True)


class TestPresentationSaveAndLoad:
    """Test saving and loading presentations with duplicated slides."""

    def test_duplicate_slides_save_and_load(self, presentation_with_slides, app):
        """Test that duplicated slides persist after save and load."""
        # Get original slide content
        original_cache_0 = presentation_with_slides._fetch_slide_cache(0)
        original_text_0 = " ".join(original_cache_0.slide_text)

        original_cache_1 = presentation_with_slides._fetch_slide_cache(1)
        original_text_1 = " ".join(original_cache_1.slide_text)

        # Duplicate slide 0 and insert at position 2
        # Note: When insert_location > source_location, the insert_location is adjusted
        initial_count = presentation_with_slides.slide_count()
        added = presentation_with_slides.duplicate_slides(0, insert_location=2)
        assert added == 1
        assert presentation_with_slides.slide_count() == initial_count + 1

        # Find where the duplicate actually ended up by checking all slides
        duplicate_found = False
        duplicate_index = None
        for i in range(presentation_with_slides.slide_count()):
            cache = presentation_with_slides._fetch_slide_cache(i)
            text = " ".join(cache.slide_text)
            if text == original_text_0 and i != 0:  # Found duplicate (not the original)
                duplicate_found = True
                duplicate_index = i
                break

        assert duplicate_found, "Duplicated slide not found"

        # Save the presentation
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            presentation_with_slides.saveas(tmp_path)
            assert Path(tmp_path).exists()

            # Load the presentation back
            loaded_prs = app.open_presentation(tmp_path)
            assert loaded_prs is not None
            assert loaded_prs.slide_count() == initial_count + 1

            # Verify the duplicated slide is present and has correct content
            duplicated_cache = loaded_prs._fetch_slide_cache(duplicate_index)
            duplicated_text = " ".join(duplicated_cache.slide_text)
            assert original_text_0 == duplicated_text, (
                f"Duplicated slide content mismatch: expected '{original_text_0}', got '{duplicated_text}'"
            )

            # Verify original slides are still intact
            loaded_cache_0 = loaded_prs._fetch_slide_cache(0)
            loaded_text_0 = " ".join(loaded_cache_0.slide_text)
            assert original_text_0 == loaded_text_0

            loaded_cache_1 = loaded_prs._fetch_slide_cache(1)
            loaded_text_1 = " ".join(loaded_cache_1.slide_text)
            assert original_text_1 == loaded_text_1

        finally:
            Path(tmp_path).unlink(missing_ok=True)

    def test_duplicate_multiple_slides_save_and_load(self, app):
        """Test that multiple duplicated slides persist after save and load."""
        # Create a presentation with multiple slides
        prs = app.new_presentation()
        if prs.slide_count() > 0:
            prs.delete_slide(0)

        # Add slides with unique content
        for i in range(3):
            prs.insert_blank_slide(i)
            slide = prs.prs.slides[i]
            textbox = slide.shapes.add_textbox(0, 0, 100, 100)
            textbox.text_frame.text = f"Original Slide {i}"

        initial_count = prs.slide_count()

        # Get original content counts before duplication
        original_content_counts = {}
        for i in range(initial_count):
            cache = prs._fetch_slide_cache(i)
            text = " ".join(cache.slide_text)
            original_content_counts[text] = original_content_counts.get(text, 0) + 1

        # Duplicate slides 0 and 2
        added = prs.duplicate_slides([0, 2], insert_location=1)
        assert added == 2
        assert prs.slide_count() == initial_count + 2

        # Verify content counts after duplication (before save)
        content_counts_before = {}
        for i in range(prs.slide_count()):
            cache = prs._fetch_slide_cache(i)
            text = " ".join(cache.slide_text)
            content_counts_before[text] = content_counts_before.get(text, 0) + 1

        # Original Slide 0 should appear twice (original + duplicate)
        assert content_counts_before.get("Original Slide 0", 0) == 2
        # Original Slide 1 should appear once (original only)
        assert content_counts_before.get("Original Slide 1", 0) == 1
        # Original Slide 2 should appear twice (original + duplicate)
        assert content_counts_before.get("Original Slide 2", 0) == 2

        # Save the presentation
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            prs.saveas(tmp_path)
            assert Path(tmp_path).exists()

            # Load the presentation back
            loaded_prs = app.open_presentation(tmp_path)
            assert loaded_prs is not None
            assert loaded_prs.slide_count() == initial_count + 2

            # Verify content counts after load
            content_counts_after = {}
            for i in range(loaded_prs.slide_count()):
                cache = loaded_prs._fetch_slide_cache(i)
                text = " ".join(cache.slide_text)
                content_counts_after[text] = content_counts_after.get(text, 0) + 1

            # Verify all content is preserved correctly
            assert content_counts_after.get("Original Slide 0", 0) == 2, (
                f"Expected 2 copies of 'Original Slide 0' after load, found {content_counts_after.get('Original Slide 0', 0)}"
            )
            assert content_counts_after.get("Original Slide 1", 0) == 1, (
                f"Expected 1 copy of 'Original Slide 1' after load, found {content_counts_after.get('Original Slide 1', 0)}"
            )
            assert content_counts_after.get("Original Slide 2", 0) == 2, (
                f"Expected 2 copies of 'Original Slide 2' after load, found {content_counts_after.get('Original Slide 2', 0)}"
            )

            # Verify content counts match before and after save/load
            assert content_counts_before == content_counts_after, (
                f"Content counts changed after save/load: before={content_counts_before}, after={content_counts_after}"
            )

        finally:
            Path(tmp_path).unlink(missing_ok=True)

    def test_duplicate_replace_save_and_load(self, presentation_with_slides, app):
        """Test duplicate, replace text, save, and load workflow."""
        # Get original slide 0 content
        original_cache_0 = presentation_with_slides._fetch_slide_cache(0)
        original_text_0 = " ".join(original_cache_0.slide_text)

        # Duplicate slide 0
        initial_count = presentation_with_slides.slide_count()
        presentation_with_slides.duplicate_slides(0, insert_location=2)

        # Find where the duplicate actually ended up
        duplicate_index = None
        for i in range(presentation_with_slides.slide_count()):
            cache = presentation_with_slides._fetch_slide_cache(i)
            text = " ".join(cache.slide_text)
            if text == original_text_0 and i != 0:  # Found duplicate
                duplicate_index = i
                break

        assert duplicate_index is not None, "Duplicated slide not found"

        # Replace text in the duplicated slide
        presentation_with_slides._replace_texts_in_slide_shapes(duplicate_index, {"Slide 0": "Replaced Text"})

        # Verify replacement worked before saving
        replaced_cache_before = presentation_with_slides._fetch_slide_cache(duplicate_index)
        replaced_text_before = " ".join(replaced_cache_before.slide_text)
        assert "Replaced Text" in replaced_text_before, f"Replacement failed before save: {replaced_text_before}"

        # Save the presentation
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            presentation_with_slides.saveas(tmp_path)

            # Load the presentation back
            loaded_prs = app.open_presentation(tmp_path)
            assert loaded_prs is not None
            assert loaded_prs.slide_count() == initial_count + 1

            # Verify the replaced text is preserved in the duplicated slide
            replaced_cache = loaded_prs._fetch_slide_cache(duplicate_index)
            replaced_text = " ".join(replaced_cache.slide_text)
            assert "Replaced Text" in replaced_text, f"Replaced text not found in slide {duplicate_index}: {replaced_text}"

            # Verify original slide 0 is unchanged
            original_cache = loaded_prs._fetch_slide_cache(0)
            original_text = " ".join(original_cache.slide_text)
            assert "Slide 0" in original_text, f"Original slide should still contain 'Slide 0': {original_text}"
            assert "Replaced Text" not in original_text, f"Original slide should not contain replaced text: {original_text}"

        finally:
            Path(tmp_path).unlink(missing_ok=True)


class TestPresentationEdgeCases:
    """Test edge cases and error handling."""

    def test_slide_index_to_ID_invalid_index(self, empty_presentation):
        """Test slide_index_to_ID() with invalid index."""
        with pytest.raises(IndexError):
            empty_presentation.slide_index_to_ID(999)

    def test_slide_ID_to_index_invalid_id(self, empty_presentation):
        """Test slide_ID_to_index() with invalid ID."""
        result = empty_presentation.slide_ID_to_index(99999)
        assert result == -1  # find_slide_by_id returns -1 for not found

    def test_delete_slide_invalid_index(self, empty_presentation):
        """Test delete_slide() with invalid index."""
        initial_count = empty_presentation.slide_count()
        with pytest.raises(IndexError):
            empty_presentation.delete_slide(999)
        # Count should not change
        assert empty_presentation.slide_count() == initial_count

    def test_duplicate_slides_empty_list(self, presentation_with_slides):
        """Test duplicate_slides() with empty list."""
        initial_count = presentation_with_slides.slide_count()
        added = presentation_with_slides.duplicate_slides([], insert_location=0)
        assert added == 0
        assert presentation_with_slides.slide_count() == initial_count

    def test_duplicate_slides_single_item_list(self, presentation_with_slides):
        """Test duplicate_slides() treats single-item list as int."""
        initial_count = presentation_with_slides.slide_count()
        added = presentation_with_slides.duplicate_slides([0], insert_location=1)
        assert added == 1
        assert presentation_with_slides.slide_count() == initial_count + 1
